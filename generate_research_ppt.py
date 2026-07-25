#!/usr/bin/env python3
"""Generate PowerPoint from the Lattice EDM circularity research paper content."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
FIGS = ROOT / "paper_figures"
OUTPUT = ROOT / "Lattice_EDM_Circularity_Research_Presentation.pptx"

# Palette
NAVY = RGBColor(0x1A, 0x2B, 0x3C)
TEAL = RGBColor(0x1F, 0x6F, 0x8B)
ACCENT = RGBColor(0xC4, 0x5C, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF7, 0xF9)
DARK = RGBColor(0x2C, 0x3E, 0x50)
MUTED = RGBColor(0x55, 0x65, 0x72)


def set_run(run, size=18, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color=LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_title_bar(slide, title: str):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.95), NAVY)
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.22), Inches(12.5), Inches(0.55))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run(run, 26, True, WHITE, "Calibri")


def add_footer(slide, page: int, total: int):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Lattice EDM Circularity — Research Presentation   |   {page}/{total}"
    set_run(run, 11, False, MUTED)
    p.alignment = PP_ALIGN.RIGHT


def bullets(slide, left, top, width, height, items, size=17):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "•  " + item
        set_run(run, size, False, DARK)


def add_picture_safe(slide, path: Path, left, top, width=None, height=None):
    if not path.exists():
        return None
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    return slide.shapes.add_picture(str(path), left, top, **kwargs)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slides_meta = []

    def new_slide(title=None):
        s = prs.slides.add_slide(blank)
        add_bg(s, LIGHT)
        if title:
            add_title_bar(s, title)
        slides_meta.append(s)
        return s

    # 1 Title
    s = new_slide()
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    add_rect(s, Inches(0), Inches(5.9), Inches(13.333), Inches(1.6), TEAL)
    box = s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Physics-Informed Machine Learning for Predicting\n"
        "Supporting-Boundary Circularity at Any Tool Position\n"
        "in Micro-EDM of Metallic Lattices"
    )
    set_run(run, 30, True, WHITE)
    p.alignment = PP_ALIGN.LEFT
    box2 = s.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11), Inches(1.0))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "Research Presentation  ·  SEM Validation  ·  Gradient Boosting  ·  Gaussian Processes"
    set_run(r2, 16, False, RGBColor(0xD0, 0xE4, 0xEF))
    box3 = s.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11), Inches(0.8))
    tf3 = box3.text_frame
    p3 = tf3.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "Lattice Circularity Analyzer Project Team"
    set_run(r3, 18, True, WHITE)

    # 2 Agenda
    s = new_slide("Agenda")
    bullets(s, Inches(0.7), Inches(1.3), Inches(11.5), Inches(5.5), [
        "Problem: oversized EDM tool on a fine metallic lattice",
        "Experimental evidence from 16 SEM-labeled runs (graphs)",
        "Ideation: circularity = process intensity × local geometry",
        "How 16 experiments become a trainable dataset",
        "Machine learning in detail — features, models, train/validate",
        "Predicting circularity ratio for any (I, T, D, x, y)",
        "Final recommended parameters and conclusions",
    ], size=20)

    # 3 Problem
    s = new_slide("The Problem")
    bullets(s, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.5), [
        "Unit cell = 500 µm; pore/node Ø ≈ 235.6 µm",
        "EDM tool tip = 900 µm → tool/pore ratio = 3.82",
        "Tool always hits pores + nodes + supporting struts together",
        "Goal: continuous nearly circular supporting ring",
        "Nodes may be destroyed — that is acceptable",
        "Hole-deviation numbers alone mislead (Run 5 paradox)",
        "Need: predict circularity ratio at any landing (x, y)",
    ], size=17)
    add_picture_safe(s, FIGS / "fig_geometry_ideation.png", Inches(7.0), Inches(1.3), width=Inches(5.7))

    # 4 SEM evidence
    s = new_slide("SEM Ground Truth — 16 Experiments")
    add_picture_safe(s, ROOT / "ACTUAL IMAGE OF THE 16 DATASETS .png", Inches(0.5), Inches(1.15), width=Inches(8.0))
    bullets(s, Inches(8.7), Inches(1.3), Inches(4.2), Inches(5.5), [
        "Only Run 4 PASS: 4 A, 150 µs, 80%",
        "15 / 16 runs destroy or irregularize the ring",
        "SEM is the true success metric",
        "Deviation ranking ≠ SEM ranking",
    ], size=16)

    # 5 Paradox graphs
    s = new_slide("The Central Paradox (Graphs)")
    add_picture_safe(s, FIGS / "fig_deviation_vs_circularity.png", Inches(0.3), Inches(1.15), width=Inches(6.4))
    add_picture_safe(s, FIGS / "fig_dev_score_ranking.png", Inches(6.9), Inches(1.15), width=Inches(6.0))
    box = s.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12), Inches(0.6))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Run 5: best deviation score → FAIL SEM   |   Run 4: near-worst deviation → ONLY PASS"
    set_run(r, 16, True, ACCENT)

    # 6 More experiment graphs
    s = new_slide("Experimental Graphs — Deviation & Asymmetry")
    add_picture_safe(s, FIGS / "fig_top_bottom_deviation.png", Inches(0.3), Inches(1.15), width=Inches(6.4))
    add_picture_safe(s, FIGS / "fig_asymmetry.png", Inches(6.9), Inches(1.15), width=Inches(6.0))

    # 7 Parameter effects
    s = new_slide("Parameter Effects on Deviation")
    add_picture_safe(s, FIGS / "fig_param_effects_dev.png", Inches(0.3), Inches(1.15), width=Inches(6.4))
    add_picture_safe(s, FIGS / "fig_duty_effects_dev.png", Inches(6.9), Inches(1.15), width=Inches(6.0))

    # 8 Circularity outcomes
    s = new_slide("SEM Circularity Outcomes")
    add_picture_safe(s, FIGS / "fig_circularity_by_run.png", Inches(0.3), Inches(1.15), width=Inches(6.5))
    add_picture_safe(s, FIGS / "fig_pass_fail_pie.png", Inches(7.1), Inches(1.2), width=Inches(5.5))

    # 9 Energy / heatmap
    s = new_slide("Energy & Current × Pulse-on Patterns")
    add_picture_safe(s, FIGS / "fig_energy_vs_circularity.png", Inches(0.3), Inches(1.15), width=Inches(6.4))
    add_picture_safe(s, FIGS / "fig_current_pulse_heatmap.png", Inches(6.9), Inches(1.15), width=Inches(6.0))

    # 10 Ideation
    s = new_slide("Ideation — The Working Mind of the Solution")
    bullets(s, Inches(0.5), Inches(1.25), Inches(12.2), Inches(5.5), [
        "Idea 1: Train on SEM boundary truth — not hole-deviation proxies",
        "Idea 2: Geometry is a first-class feature (strut/node overlap under the tool)",
        "Idea 3: Expand sparse lab truth with physics-structured data (GP + spatial replay)",
        "Circularity = joint effect of EDM intensity AND local lattice geometry at (x, y)",
        "Same recipe can PASS at pore center and FAIL near a strut",
    ], size=19)

    # 11 Pipeline
    s = new_slide("Prediction Pipeline")
    add_picture_safe(s, FIGS / "fig_ml_pipeline.png", Inches(0.6), Inches(1.2), width=Inches(12.0))

    # 12 Data expansion
    s = new_slide("From 16 Experiments → Trainable Data")
    add_picture_safe(s, FIGS / "fig_data_expansion.png", Inches(0.5), Inches(1.15), width=Inches(12.2))
    bullets(s, Inches(0.7), Inches(4.3), Inches(12), Inches(2.5), [
        "Augmentation A: Gaussian Process posterior → 1,100 synthetic EDM points (not new lab trials)",
        "Augmentation B: 16 SEM labels × ~25 landings → ~400 position-aware rows",
        "Label rule: c_train = clip(c_SEM − 2.5·geometry_risk + edm_bonus, 1, 5)",
    ], size=16)

    # 13 Features
    s = new_slide("ML Features — 20-D Vector")
    bullets(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.5), [
        "EDM block (7): I, T, D, Energy E=I·T·D/100, pulse-off proxy, I·D, T/D",
        "Geometry block (13): x/W, y/W, dist to strut/node, nodes/pores in tool, strut intersection, overlaps, geometry risk, tool/pore ratio, working area, tool Ø",
        "Why: trees can split on intensity AND local damage exposure",
        "Same (I,T,D) with different (x,y) → different prediction",
    ], size=16)
    bullets(s, Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.5), [
        "Targets:",
        "A) Circularity score c ∈ [1,5] from SEM",
        "B) Supporting intact s ∈ {0,1}",
        "Reported ratio = c / 5",
        "PASS if c ≥ 3.5, support OK, risk ≤ 0.55",
        "NOT trained to minimize Hole_Dev",
    ], size=16)

    # 14 Models
    s = new_slide("Models, Training & Validation")
    bullets(s, Inches(0.5), Inches(1.2), Inches(12.2), Inches(5.8), [
        "Gradient Boosting (circularity): 120 trees, max depth 4 → score 1–5",
        "Gradient Boosting (supporting): 80 trees, max depth 3 → threshold 0.5",
        "Gaussian Process (Matérn 5/2 + WhiteKernel) on 16 SEM labels to search robust (I,T,D)",
        "Polynomial Ridge LOOCV on 16 runs → MAE ≈ 1.02 / 5 (small-n honesty check)",
        "Train: load 16 runs → expand on landing grid → fit both GBRs → save joblib bundle",
        "Validate: LOOCV + consistency tests (gentle recipes peak at pore centers; I≥8 stays FAIL)",
        "Inference blend: c = (1−w)·c_ML + w·c_H, where w rises if tool/pore drifts from lab geometry",
    ], size=16)

    # 15 Any position
    s = new_slide("Predict Circularity Ratio for Any Position")
    bullets(s, Inches(0.6), Inches(1.25), Inches(12), Inches(5.5), [
        "1. Geometry pass at (x, y) → risk, distances, overlaps",
        "2. Build 20-D feature row (EDM + geometry)",
        "3. ML forward: circularity GBR + supporting GBR",
        "4. Heuristic forward: gentle/aggressive rules + geometry risk",
        "5. Blend with drift weight w",
        "6. Report score, ratio = score/5, supporting OK, PASS/FAIL",
        "Repeat on a grid → spatial circularity field for fixed EDM settings",
    ], size=18)

    # 16 Final recommendations
    s = new_slide("Final Recommended Parameters")
    add_picture_safe(s, FIGS / "fig_final_recommendations.png", Inches(0.4), Inches(1.15), width=Inches(7.0))
    bullets(s, Inches(7.6), Inches(1.3), Inches(5.2), Inches(5.5), [
        "Unknown position (robust): 4 A, 150 µs, 80%",
        "Pore center: 4 A, 150 µs, 80%",
        "Mid pore: 4 A, 148 µs, 79%",
        "Near strut: 3.5 A, 150 µs, 78%",
        "Near node: 3.5 A, 145 µs, 76%",
        "Reject deviation-only: 6 A, 50 µs, 64%",
        "Pattern: low I + long T + high D",
    ], size=15)

    # 17 Conclusion
    s = new_slide("Conclusion")
    bullets(s, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5), [
        "SEM — not hole deviation — defines success for oversized-tool lattice EDM",
        "Winning pattern: low current + long pulse-on + high duty (Run 4)",
        "16 labels become useful via GP synthesis + geometry-risk spatial replay",
        "20-D EDM+geometry features + Gradient Boosting + physics blend enable any-position prediction",
        "Circularity ratio = score/5 is computable for any (I, T, D, x, y)",
        "Next: more SEM trials in the 3.5–4.5 A / 140–150 µs / 76–80% island",
    ], size=17)

    # 18 Thank you
    s = new_slide()
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    box = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Thank You"
    set_run(r, 48, True, WHITE)
    p.alignment = PP_ALIGN.CENTER
    box2 = s.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(1.2))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "Questions & Discussion"
    set_run(r2, 22, False, RGBColor(0xD0, 0xE4, 0xEF))
    p2.alignment = PP_ALIGN.CENTER

    total = len(slides_meta)
    # footers on content slides (skip title and thank you: index 0 and -1)
    for i, slide in enumerate(slides_meta):
        if i == 0 or i == total - 1:
            continue
        add_footer(slide, i + 1, total)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT} ({total} slides)")
    return total


if __name__ == "__main__":
    build()
