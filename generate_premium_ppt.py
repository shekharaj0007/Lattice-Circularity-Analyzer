#!/usr/bin/env python3
"""
Premium 18–20 slide PowerPoint for Lattice EDM Circularity project.
Includes research graphs, SEM imagery, and website UI illustrations.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parent
FIGS = ROOT / "paper_figures"
ASSETS = ROOT / "assets"
OUTPUT = ROOT / "Lattice_EDM_Circularity_Premium_Presentation.pptx"

# Design system — engineering / research (not purple-AI default)
NAVY = RGBColor(0x0F, 0x1C, 0x2E)
INK = RGBColor(0x1B, 0x2A, 0x3A)
TEAL = RGBColor(0x0E, 0x7C, 0x86)
AMBER = RGBColor(0xD9, 0x7A, 0x0C)
CREAM = RGBColor(0xF7, 0xF4, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x3D, 0x4F, 0x5F)
SOFT = RGBColor(0xE8, 0xEE, 0xF2)
GREEN = RGBColor(0x1B, 0x7A, 0x3D)
RED = RGBColor(0xB3, 0x3A, 0x3A)


def font(run, size=18, bold=False, color=INK, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    return sh


def round_rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    try:
        sh.adjustments[0] = 0.08
    except Exception:
        pass
    return sh


def bg(slide, color=CREAM):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def header(slide, title: str, subtitle: str | None = None):
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.05), NAVY)
    rect(slide, Inches(0), Inches(1.05), Inches(13.333), Inches(0.08), TEAL)
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(12.4), Inches(0.55))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    font(r, 26, True, WHITE)
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.45), Inches(0.68), Inches(12.4), Inches(0.3))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        font(r2, 12, False, RGBColor(0xB8, 0xC7, 0xD4))


def footer(slide, n, total):
    rect(slide, Inches(0), Inches(7.15), Inches(13.333), Inches(0.35), SOFT)
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.18), Inches(8), Inches(0.28))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "LatticeFlow  ·  EDM Circularity Research"
    font(r, 10, False, SLATE)
    box2 = slide.shapes.add_textbox(Inches(10.5), Inches(7.18), Inches(2.4), Inches(0.28))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{n}  /  {total}"
    font(r2, 10, True, SLATE)


def textbox(slide, l, t, w, h, text, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    font(r, size, bold, color)
    return box


def bullets(slide, l, t, w, h, items, size=16, color=INK):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        r = p.add_run()
        r.text = "▸  " + item
        font(r, size, False, color)
    return box


def pic(slide, path: Path, l, t, w=None, h=None, width=None, height=None):
    if not Path(path).exists():
        return None
    w = width if width is not None else w
    h = height if height is not None else h
    kw = {}
    if w is not None:
        kw["width"] = w
    if h is not None:
        kw["height"] = h
    return slide.shapes.add_picture(str(path), l, t, **kw)


def card(slide, l, t, w, h, title, body_lines, accent=TEAL):
    round_rect(slide, l, t, w, h, WHITE)
    rect(slide, l, t, Inches(0.12), h, accent)
    textbox(slide, l + Inches(0.28), t + Inches(0.18), w - Inches(0.4), Inches(0.35), title, 15, True, NAVY)
    bullets(slide, l + Inches(0.22), t + Inches(0.55), w - Inches(0.35), h - Inches(0.65), body_lines, 13, SLATE)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    all_slides = []

    def add(title=None, subtitle=None, cream=True):
        s = prs.slides.add_slide(blank)
        bg(s, CREAM if cream else NAVY)
        if title:
            header(s, title, subtitle)
        all_slides.append(s)
        return s

    # ========== 1 TITLE ==========
    s = add(cream=False)
    rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    rect(s, Inches(0), Inches(0), Inches(0.18), Inches(7.5), TEAL)
    rect(s, Inches(0), Inches(6.35), Inches(13.333), Inches(1.15), RGBColor(0x14, 0x28, 0x3E))
    textbox(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.4),
            "RESEARCH  ·  MACHINE LEARNING  ·  MICRO-EDM", 13, True, TEAL)
    textbox(s, Inches(0.7), Inches(2.0), Inches(12), Inches(2.2),
            "Predicting Supporting-Boundary Circularity\nat Any Tool Position in Lattice EDM",
            34, True, WHITE)
    textbox(s, Inches(0.7), Inches(4.5), Inches(11.5), Inches(0.8),
            "Physics-informed Gradient Boosting  ·  Gaussian Process data expansion  ·  SEM-validated outcomes\n"
            "Interactive LatticeFlow analyzer for spatial circularity prediction",
            15, False, RGBColor(0xB8, 0xC7, 0xD4))
    textbox(s, Inches(0.7), Inches(6.6), Inches(10), Inches(0.5),
            "Lattice Circularity Analyzer Project  ·  Premium Presentation Deck", 14, True, WHITE)

    # ========== 2 AGENDA ==========
    s = add("Agenda", "What this presentation covers")
    items = [
        ("01", "Problem & Geometry", "Oversized tool on a fine lattice"),
        ("02", "SEM Evidence & Graphs", "16 runs, paradox, parameter effects"),
        ("03", "Ideation & ML Pipeline", "How prediction becomes possible"),
        ("04", "Train · Validate · Infer", "Features, models, any-position path"),
        ("05", "LatticeFlow Website", "Live analyzer UI walkthrough"),
        ("06", "Final Answers", "Recommended EDM settings by zone"),
    ]
    for i, (num, title, sub) in enumerate(items):
        col = i % 3
        row = i // 3
        l = Inches(0.45 + col * 4.2)
        t = Inches(1.5 + row * 2.5)
        round_rect(s, l, t, Inches(3.95), Inches(2.15), WHITE)
        rect(s, l, t, Inches(3.95), Inches(0.12), TEAL if row == 0 else AMBER)
        textbox(s, l + Inches(0.25), t + Inches(0.35), Inches(3.4), Inches(0.4), num, 22, True, TEAL)
        textbox(s, l + Inches(0.25), t + Inches(0.85), Inches(3.4), Inches(0.4), title, 18, True, NAVY)
        textbox(s, l + Inches(0.25), t + Inches(1.35), Inches(3.4), Inches(0.5), sub, 13, False, SLATE)

    # ========== 3 PROBLEM ==========
    s = add("The Engineering Problem", "900 µm tool vs 235.6 µm pore — ratio 3.82×")
    card(s, Inches(0.4), Inches(1.4), Inches(5.9), Inches(5.3), "What must succeed", [
        "Unit cell side = 500 µm",
        "Pore / node diameter ≈ 235.6 µm",
        "Tool tip diameter = 900 µm",
        "Tool overlaps pore + nodes + struts at once",
        "Supporting ring (black) must stay continuous & circular",
        "Nodes (red) may be destroyed — acceptable",
        "Need circularity ratio for ANY landing (x, y)",
    ])
    pic(s, FIGS / "fig_geometry_ideation.png", Inches(6.55), Inches(1.45), width=Inches(6.3))

    # ========== 4 SEM ==========
    s = add("SEM Ground Truth", "Only visual truth of supporting-boundary circularity")
    pic(s, ROOT / "ACTUAL IMAGE OF THE 16 DATASETS .png", Inches(0.35), Inches(1.25), width=Inches(8.3))
    round_rect(s, Inches(8.9), Inches(1.4), Inches(4.0), Inches(5.2), WHITE)
    rect(s, Inches(8.9), Inches(1.4), Inches(4.0), Inches(0.12), GREEN)
    textbox(s, Inches(9.1), Inches(1.7), Inches(3.6), Inches(0.4), "Key finding", 16, True, NAVY)
    bullets(s, Inches(9.05), Inches(2.25), Inches(3.7), Inches(4.0), [
        "16 lab EDM trials imaged by SEM",
        "ONLY Run 4 PASS",
        "4 A · 150 µs · 80%",
        "15 / 16 fail the ring test",
        "SEM > hole-deviation metrics",
        "ML must maximize SEM circularity",
    ], 14)

    # ========== 5 PARADOX ==========
    s = add("The Central Paradox", "Best numerical deviation ≠ best supporting ring")
    pic(s, FIGS / "fig_deviation_vs_circularity.png", Inches(0.3), Inches(1.25), width=Inches(6.4))
    pic(s, FIGS / "fig_dev_score_ranking.png", Inches(6.85), Inches(1.25), width=Inches(6.1))
    round_rect(s, Inches(0.4), Inches(6.35), Inches(12.5), Inches(0.65), WHITE)
    textbox(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.45),
            "Run 5 (6 A, 50 µs, 64%) wins deviation ranking but FAILS SEM   ·   "
            "Run 4 (4 A, 150 µs, 80%) is near-worst by deviation yet is the ONLY PASS",
            13, True, ACCENT if False else AMBER)

    # ========== 6 DEVIATION GRAPHS ==========
    s = add("Experimental Graphs — Top/Bottom Deviation & Asymmetry")
    pic(s, FIGS / "fig_top_bottom_deviation.png", Inches(0.25), Inches(1.25), width=Inches(6.4))
    pic(s, FIGS / "fig_asymmetry.png", Inches(6.8), Inches(1.25), width=Inches(6.15))

    # ========== 7 PARAMETER EFFECTS ==========
    s = add("Parameter Effects", "How current, pulse-on, and duty shape deviation")
    pic(s, FIGS / "fig_param_effects_dev.png", Inches(0.25), Inches(1.25), width=Inches(6.4))
    pic(s, FIGS / "fig_duty_effects_dev.png", Inches(6.8), Inches(1.25), width=Inches(6.15))

    # ========== 8 CIRCULARITY OUTCOMES ==========
    s = add("SEM Circularity Outcomes", "Pass/fail distribution across the DOE")
    pic(s, FIGS / "fig_circularity_by_run.png", Inches(0.25), Inches(1.25), width=Inches(6.5))
    pic(s, FIGS / "fig_pass_fail_pie.png", Inches(7.0), Inches(1.35), width=Inches(5.7))

    # ========== 9 ENERGY / HEATMAP / MRR ==========
    s = add("Energy, Heatmap & Process Responses")
    pic(s, FIGS / "fig_energy_vs_circularity.png", Inches(0.2), Inches(1.2), width=Inches(4.3))
    pic(s, FIGS / "fig_current_pulse_heatmap.png", Inches(4.55), Inches(1.2), width=Inches(4.2))
    pic(s, FIGS / "fig_mrr_twr.png", Inches(8.85), Inches(1.2), width=Inches(4.2))

    # ========== 10 IDEATION ==========
    s = add("Ideation — Working Mind of the Solution", "Circularity = process intensity × local geometry")
    card(s, Inches(0.4), Inches(1.4), Inches(4.0), Inches(5.2), "Idea 1 — SEM objective", [
        "Train on supporting-ring circularity",
        "Do NOT minimize Hole_Dev alone",
        "Run 5 is the wrong optimum",
        "Run 4 is the reference success",
    ], TEAL)
    card(s, Inches(4.65), Inches(1.4), Inches(4.0), Inches(5.2), "Idea 2 — Geometry features", [
        "Tool footprint vs struts/nodes",
        "Distances, overlaps, intersection",
        "Geometry risk index ∈ [0,1]",
        "Same recipe, different (x,y) outcome",
    ], AMBER)
    card(s, Inches(8.9), Inches(1.4), Inches(4.0), Inches(5.2), "Idea 3 — Sparse-data expansion", [
        "GP posterior → 1,100 EDM synthetics",
        "16 × landing grid → ~400 rows",
        "Risk-modulated SEM label replay",
        "Physics blend at inference time",
    ], GREEN)

    # ========== 11 PIPELINE ==========
    s = add("End-to-End Prediction Pipeline")
    pic(s, FIGS / "fig_ml_pipeline.png", Inches(0.45), Inches(1.25), width=Inches(12.4))

    # ========== 12 DATA EXPANSION ==========
    s = add("From 16 Experiments to Trainable Data")
    pic(s, FIGS / "fig_data_expansion.png", Inches(0.5), Inches(1.25), width=Inches(12.2))
    round_rect(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.3), WHITE)
    bullets(s, Inches(0.75), Inches(4.7), Inches(11.8), Inches(1.9), [
        "Augmentation A: Gaussian Process posterior sampling densifies (I, T, D) space → 1,100 synthetics (tagged GP_posterior; not new lab trials)",
        "Augmentation B: each SEM label replayed on ~25 landings → c_train = clip(c_SEM − 2.5·geometry_risk + edm_bonus, 1, 5)",
        "Result: models learn that gentle recipes degrade when the oversized tool sits on dense strut intersections",
    ], 15)

    # ========== 13 ML DETAIL ==========
    s = add("Machine Learning Detail", "Features · models · train · validate · blend")
    card(s, Inches(0.35), Inches(1.35), Inches(4.15), Inches(5.4), "20-D Feature Vector", [
        "EDM (7): I, T, D, Energy, pulse-off, I·D, T/D",
        "Geometry (13): x/W, y/W, distances,",
        "nodes/pores in tool, strut length,",
        "overlaps, geometry risk, ratio, W, tool Ø",
        "Targets: score 1–5 + support 0/1",
        "Ratio reported as score / 5",
    ], TEAL)
    card(s, Inches(4.65), Inches(1.35), Inches(4.15), Inches(5.4), "Algorithms", [
        "GBR circularity: 120 trees, depth 4",
        "GBR supporting: 80 trees, depth 3",
        "GP Matérn 5/2 + WhiteKernel on 16 SEM",
        "Ridge poly LOOCV → MAE ≈ 1.02 / 5",
        "Shallow trees limit small-n overfitting",
        "GP ranks μ + 0.15·σ candidates",
    ], AMBER)
    card(s, Inches(8.95), Inches(1.35), Inches(4.0), Inches(5.4), "Inference Blend", [
        "c = (1−w)·c_ML + w·c_H",
        "w rises if tool/pore drifts from lab",
        "Heuristic encodes gentle vs blast rules",
        "PASS: c≥3.5, support OK, risk≤0.55",
        "Consistency tests replace huge holdout",
        "SEM remains ultimate acceptance",
    ], GREEN)

    # ========== 14 ANY POSITION ==========
    s = add("Predict Circularity Ratio for Any Position", "Runtime path for query (I, T, D, x, y)")
    steps = [
        ("1", "Geometry pass", "analyze_position(x,y) → risk & overlaps"),
        ("2", "Feature build", "Concatenate 7 EDM + 13 geometry = 20-D"),
        ("3", "ML forward", "GBR score + supporting probability"),
        ("4", "Heuristic", "Gentle/aggressive rules + risk penalty"),
        ("5", "Blend", "Drift-weighted mix of ML and physics"),
        ("6", "Report", "Score, ratio=score/5, PASS/FAIL"),
    ]
    for i, (n, title, desc) in enumerate(steps):
        col = i % 3
        row = i // 3
        l = Inches(0.4 + col * 4.25)
        t = Inches(1.45 + row * 2.6)
        round_rect(s, l, t, Inches(4.05), Inches(2.3), WHITE)
        oval = s.shapes.add_shape(MSO_SHAPE.OVAL, l + Inches(0.25), t + Inches(0.35), Inches(0.55), Inches(0.55))
        oval.fill.solid()
        oval.fill.fore_color.rgb = TEAL
        oval.line.fill.background()
        textbox(s, l + Inches(0.25), t + Inches(0.42), Inches(0.55), Inches(0.4), n, 16, True, WHITE, PP_ALIGN.CENTER)
        textbox(s, l + Inches(1.0), t + Inches(0.4), Inches(2.8), Inches(0.4), title, 18, True, NAVY)
        textbox(s, l + Inches(0.3), t + Inches(1.15), Inches(3.5), Inches(0.8), desc, 14, False, SLATE)

    # ========== 15 SITE INTRO ==========
    s = add("LatticeFlow Website — Overview", "Interactive circularity analyzer deployed for engineers")
    bullets(s, Inches(0.5), Inches(1.35), Inches(5.8), Inches(5.3), [
        "Enter EDM parameters: I, T, duty",
        "Set tool / pore / working area geometry",
        "Choose landing position (x, y)",
        "Get circularity score, ratio, PASS/FAIL",
        "Run full-grid heatmap scans",
        "AI recommends best landing position",
        "Auto engineering report + chat assistant",
        "Supports additional tool shapes",
    ], 16)
    pic(s, ASSETS / "Grid Scan For Circularity.png", Inches(6.5), Inches(1.35), width=Inches(6.4))

    # ========== 16 SITE GRID / AI ==========
    s = add("Website — Grid Scan & AI Best Position")
    pic(s, ASSETS / "Grid Scan Analyzation.png", Inches(0.25), Inches(1.2), width=Inches(6.35))
    pic(s, ASSETS / "Grid Scan For Best Recoomended Position By AI.png", Inches(6.75), Inches(1.2), width=Inches(6.2))

    # ========== 17 SITE PASS/FAIL ==========
    s = add("Website — PASS / FAIL Circularity Detection")
    pic(s, ASSETS / "Pass Circularity With Image Position.png", Inches(0.25), Inches(1.2), width=Inches(6.35))
    pic(s, ASSETS / "Fail Circularity With Image Position.png", Inches(6.75), Inches(1.2), width=Inches(6.2))
    round_rect(s, Inches(0.4), Inches(6.35), Inches(12.5), Inches(0.65), WHITE)
    textbox(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.45),
            "PASS gates: circularity ≥ 3.5/5  ·  ratio ≥ 0.70  ·  supporting intact  ·  geometry risk ≤ 0.55",
            14, True, GREEN)

    # ========== 18 SITE REPORT / CHAT / SHAPES ==========
    s = add("Website — Report, AI Assistant & Shape Analysis")
    pic(s, ASSETS / "Detailed Engineering Report.png", Inches(0.2), Inches(1.2), width=Inches(4.2))
    pic(s, ASSETS / "Ai Engineering Assistant.png", Inches(4.55), Inches(1.2), width=Inches(4.2))
    pic(s, ASSETS / "Analyze Other Shapes .png", Inches(8.9), Inches(1.2), width=Inches(4.15))

    # ========== 19 FINAL RECS ==========
    s = add("Final Recommended Parameters", "SEM-informed answers by landing situation")
    pic(s, FIGS / "fig_final_recommendations.png", Inches(0.3), Inches(1.2), width=Inches(7.0))
    round_rect(s, Inches(7.5), Inches(1.35), Inches(5.4), Inches(5.3), WHITE)
    rect(s, Inches(7.5), Inches(1.35), Inches(5.4), Inches(0.12), TEAL)
    textbox(s, Inches(7.75), Inches(1.65), Inches(5.0), Inches(0.4), "Recommended settings", 16, True, NAVY)
    bullets(s, Inches(7.7), Inches(2.2), Inches(5.0), Inches(4.2), [
        "Unknown position: 4 A, 150 µs, 80%",
        "Pore center: 4 A, 150 µs, 80%",
        "Mid pore: 4 A, 148 µs, 79%",
        "Near strut: 3.5 A, 150 µs, 78%",
        "Near node: 3.5 A, 145 µs, 76%",
        "Reject: 6 A, 50 µs, 64% (deviation-only)",
        "Pattern: low I + long T + high D",
        "Fine servo feed + stable flush",
    ], 14)

    # ========== 20 CONCLUSION ==========
    s = add("Conclusion & Takeaways")
    takeaways = [
        ("SEM first", "Supporting-ring circularity — not hole deviation — is the success metric."),
        ("Sparse → dense", "16 runs expand via GP synthetics + geometry-risk spatial label replay."),
        ("Any position", "20-D EDM+geometry features + GBR + physics blend → ratio at any (x,y)."),
        ("LatticeFlow", "Website turns the science into grid heatmaps, PASS/FAIL, reports, and guidance."),
    ]
    for i, (title, body) in enumerate(takeaways):
        t = Inches(1.35 + i * 1.35)
        round_rect(s, Inches(0.5), t, Inches(12.3), Inches(1.2), WHITE)
        rect(s, Inches(0.5), t, Inches(0.14), Inches(1.2), TEAL if i % 2 == 0 else AMBER)
        textbox(s, Inches(0.9), t + Inches(0.2), Inches(11.5), Inches(0.35), title, 16, True, NAVY)
        textbox(s, Inches(0.9), t + Inches(0.55), Inches(11.5), Inches(0.45), body, 14, False, SLATE)

    # ========== 21 THANK YOU ==========
    s = add(cream=False)
    rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    rect(s, Inches(0), Inches(0), Inches(0.18), Inches(7.5), AMBER)
    textbox(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.0),
            "Thank You", 48, True, WHITE, PP_ALIGN.CENTER)
    textbox(s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.6),
            "Questions & Discussion", 22, False, RGBColor(0xB8, 0xC7, 0xD4), PP_ALIGN.CENTER)
    textbox(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.8),
            "PDF research paper + this deck available in the project repository\n"
            "Lattice Circularity Analyzer  ·  LatticeFlow",
            14, False, RGBColor(0x8A, 0xA0, 0xB2), PP_ALIGN.CENTER)

    total = len(all_slides)
    for i, slide in enumerate(all_slides):
        # skip dark title/thank-you for cream footer style; still number content slides
        if i == 0 or i == total - 1:
            continue
        footer(slide, i + 1, total)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT} ({total} slides)")
    return total


if __name__ == "__main__":
    build()
